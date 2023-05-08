from pptx import Presentation
from pptx.util import Inches

class PPTGenerator:

    def __init__(self,contents):
        self.contents = contents

    def addTitleSlide(self,prs:Presentation) -> None:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.contents['Title']
        author = slide.placeholders[1]
        author.text = self.contents['Author Name']

        prs.save(self.contents["Presentation Name"])

    def addTableOfContents(self,prs:Presentation) ->None:
        


if __name__ == "__main__":

    content = { "Presentation Name" : "Sample.pptx",
                "Title":"Chapter Name",
                "Author Name":"Author Name",
                "Topics":{"Topic1":["Bullet Point T1.1","Bullet point T1.2","Bullet point T1.3"],
                        "Topic2":["Bullet Point T2.1","Bullet point T2.2","Bullet point T2.3"]
                        }
            }
    
    MyPptGenerator = PPTGenerator(contents=content)
    prs = Presentation()
    MyPptGenerator.addTitleSlide(prs)
    MyPptGenerator.addTableOfContents(prs)